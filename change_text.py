from cmd import Cmd
from dataclasses import dataclass
from datetime import date
from typing import List

from pptx import Presentation
from PIL import Image
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import os
from pptx.util import Pt

from copy import deepcopy
from pptx.table import Table, _Row, _Column, _Cell

from calculs.periodes.pr import PR1
from openpyxl import load_workbook

import shutil

def excel(Class):

    original = "database/Resultat.xlsx"
    target = "result/"+ Class.Nom + "- " + Class.Isin + ".xlsx"
    shutil.copyfile(original, target)
    #le dictionnaire qu'on va injecter dans l excel
    balises_ref = [
    {
    "Nom": "Nom du produit",
    "Balise": "<NOM>",
    "Result": Class.Nom,
    },
    {
    "Nom": "Isin",
    "Balise": "<ISIN>",
    "Result": Class.Isin,
    },
    {
    "Nom": "Droit",
    "Balise": "<droit>",
    "Result": Class.Droit,
    },
    {
    "Nom": "Emission",
    "Balise": "<émission>",
    "Result": Class.Emission,
    },
    {
    "Nom": "Commission",
    "Balise": "<COM>",
    "Result": Class.COM,
    },
    {
    "Nom": "Template",
    "Balise": "<>",
    "Result": Class.template,
    },
    {
    "Nom": "Typologie",
    "Balise": "<Typologie>",
    "Result": Class.Typologie,
    },
    {
    "Nom": "Fréquence",
    "Balise": "<F0>",
    "Result": Class.F0,
    },
     {
    "Nom": "Coupon périodique",
    "Balise": "<CPN>",
    "Result": Class.CPN,
    },
     {
    "Nom": "Sous-jacents",
    "Balise": "<NOMSOUSJACENT>",
    "Result": Class.NOMSOUSJACENT,
    },
     {
    "Nom": "Type de sous-jacent",
    "Balise": "<sous_jacent>",
    "Result": Class.sous_jacent,
    },
     {
    "Nom": "Barrière de remboursement initiale",
    "Balise": "<BAC>",
    "Result": Class.BAC,
    },
     {
    "Nom": "Type de barrière",
    "Balise": "<type_bar>",
    "Result": Class.type_bar,
    },
     {
    "Nom": "Pas de dégressivité",
    "Balise": "<DEG>",
    "Result": Class.DEG,
    },
     {
    "Nom": "Avant dernier niveau de barrière dégressive",
    "Balise": "<ABDAC>",
    "Result": Class.ABDAC,
    },
     {
    "Nom": "Dernier niveau de barrière dégressive",
    "Balise": "<DBAC>",
    "Result": Class.DBAC,
    },
     {
    "Nom": "Barrière de coupon",
    "Balise": "<BCPN>",
    "Result": Class.BCPN,
    },
     {
    "Nom": "Barrière de coupon dégressive",
    "Balise": "<BCPN_is_degressif>",
    "Result": Class.BCPN_is_degressif,
    },
       {
    "Nom": "Début dégressivité",
    "Balise": "<DDP>",
    "Result": Class.DDP,
    },
       {
    "Nom": "Barrière coupon mémoire",
    "Balise": "<CPN_is_memoire>",
    "Result": Class.CPN_is_memoire,
    },
       {
    "Nom": "Barrière de protection",
    "Balise": "<PDI>",
    "Result": Class.PDI,
    },
       {
    "Nom": "Dates de constatations initiales",
    "Balise": "<DCI>",
    "Result": Class.DCI,
    },
       {
    "Nom": "Dates de premier rappel",
    "Balise": "<DPR>",
    "Result": Class.DPR_affichage,
    },
    {
    "Nom": "Avant dernière date de constatation",
    "Balise": "<ADCF>",
    "Result": Class.ADCF_affichage,
    },
       {
    "Nom": "Date de constatation finale",
    "Balise": "<DCF>",
    "Result": Class.DCF_affichage,
    },
       {
    "Nom": "Type de strike",
    "Balise": "<type_strike>",
    "Result": Class.type_strike,
    },
       {
    "Nom": "Date de remboursement premier rappel",
    "Balise": "<DR1>",
    "Result": Class.DR1_affichage,
    },
    {
    "Nom": "Avant dernière date de remboursement",
    "Balise": "<DADR>",
    "Result": Class.DADR_affichage,
    },
    {
    "Nom": "Date d'échéance",
    "Balise": "<DEC>",
    "Result": Class.DEC_affichage,
    },
    {
    "Nom": "Jour de référence",
    "Balise": "<JDR>",
    "Result": Class.JDR,
    },
    {
    "Nom": "Nombre de jours ouvrés",
    "Balise": "<NJO>",
    "Result": Class.NJO,
    },
    {
    "Nom": "Niveau de scénario défavorable",
    "Balise": "<NSD>",
    "Result": Class.NSD,
    },
    {
    "Nom": "Niveau de scénario médian",
    "Balise": "<NSM>",
    "Result": Class.NSM,
    },
    {
    "Nom": "Niveau de scénario Final",
    "Balise": "<NSF>",
    "Result": Class.NSF,
    }]
    wb = load_workbook(target)
    ws = wb['Feuil1']
    ws_tra = wb['TRA']
    ws_date = wb['DATE']

    for i in range(len(balises_ref)):
        a = "a"+ str(i+1)
        b = "b"+ str(i+1)
        c = "c"+ str(i+1)
 
        ws[a] = balises_ref[i]["Nom"]
        ws[b] = balises_ref[i]["Balise"]
        ws[c] = balises_ref[i]["Result"]

    if  (Class.Typologie == "coupon autocall"):
        ws_date["A1"] = "Dates de constatations"
        ws_date["A2"] = Class.Datesconstatations1
        ws_date["A3"] = "Dates de remboursement"
        ws_date["A4"] = Class.Datesremb1

    else:
        ws_date["A1"] = "Dates de constatations"
        ws_date["A2"] = Class.Datesconstatations1
        ws_date["A3"] = "Dates de paiement1"
        ws_date["A4"] = Class.Datespaiement1   
        ws_date["A5"] = "Dates de remboursement"
        ws_date["A6"] = Class.Datesremb1


        ws_tra["A2"] = "1" 
        ws_tra["A3"] = "1"  
        ws_tra["A4"] = "1"  
        ws_tra["A5"] = "1"  
        ws_tra["A6"] = "1" 

    ws_tra['A1'] = "TRA.D.A"
    ws_tra['A2'] = Class.TRA_A_S1

    wb.save(target)

def elementsToReplaceDegressivite(Class, shapes):
    #si degressif on supprime les balises
    if Class.BAC_is_degressif != "":
        Class.desonndr = ""
        Class.longuephrase = ""
        Class.SDBAC = ""
        Class.PDINSM = ""
        Class.ETPDI = ""

#TOUS LES REMPLACEMENTS 
def elementsToReplaceRemplacement(Class, shapes):
    #mise en avant car transforme en bac

    replace_text({'<NDRTES>': Class.desonndr}, shapes)
    replace_text({'<DESONNDR>': Class.desonndr}, shapes)

    replace_text({'<longuephrase>': Class.longuephrase}, shapes)
    replace_text({'<SDBAC>': Class.SDBAC}, shapes)
    replace_text({'<PDINSM>': Class.PDINSM}, shapes)
    replace_text({'<TEST>': Class.ETPDI}, shapes) 

    ebac = Class.EBAC + "%"
    replace_text({'<EBAC>': ebac}, shapes) 

    replace_text({'<NOM>':  Class.Nom}, shapes)
    replace_text({'<droit>':  Class.Droit}, shapes)
    replace_text({'<ISIN>':  Class.Isin}, shapes)
    replace_text({'<émission>':  Class.Emission_affichage}, shapes)
    replace_text({'<DCI>':  Class.DCI}, shapes)
    replace_text({'<DCI_MAJ>':  Class.DCI_MAJ}, shapes)

    replace_text({'<1DR>':  Class.DR1_affichage}, shapes)
    replace_text({'<PR1+1>':  str(int(Class.PR1)+1)}, shapes)

    replace_text({'<DPR>':  Class.DPR_affichage}, shapes)
    
    replace_text({'<DCF>':  Class.DCF_affichage}, shapes)
    replace_text({'<DCF_MAJ>':  Class.DCF_MAJ}, shapes)
    replace_text({'<DEC>':  Class.DEC_affichage}, shapes)
    
    if (Class.F0 == "jours"):
        Class.F0_affichage = "jour"
    else:
        Class.F0_affichage = Class.F0

    Class.F0M = str(Class.F0).capitalize()
    Class.F0SM = str(Class.F0s).capitalize()

    replace_text({'<F0M>':  Class.F0M}, shapes)
    replace_text({'<F0SM>':  Class.F0SM}, shapes)

    replace_text({'<F0>':  Class.F0_affichage}, shapes)

    replace_text({'<DDCI>':  Class.DDCI_affichage}, shapes)

    replace_text({'<DDCI_MAJ>':  Class.DDCI_MAJ}, shapes)
    replace_text({'<DDCI_M_B_STRIKE>':  Class.DDCI_M_B_Strike}, shapes)

    replace_text({'<TSJ1>':  Class.TSJ[0]}, shapes)

    # replace_text({'<TSJ2>':  Class.TSJ2}, shapes)
    # replace_text({'<TSJ3>':  Class.TSJ3}, shapes)
    # replace_text({'<TSJ4>':  Class.TSJ4}, shapes)
    # replace_text({'<TSJ5>':  Class.TSJ5}, shapes)
    # replace_text({'<PCS1>':  Class.PCS1}, shapes)
    # replace_text({'<PCS2>':  Class.PCS2}, shapes)
    # replace_text({'<PCS3>':  Class.PCS3}, shapes)
    # replace_text({'<PCS4>':  Class.PCS4}, shapes)
    # replace_text({'<PCS5>':  Class.PCS5}, shapes)

    cpn = float(Class.CPN)
    len_cpn = (str(cpn)[::-1].find('.'))

    if (len_cpn < 2):
            cpn = (f'{float(Class.CPN):.2f}')
       
        


    cpn = str(cpn) + "%"
    cpn = cpn.replace(".", ",")
    cpn = cpn.rstrip()
    replace_text({'<CPN>':  cpn}, shapes)
    pdi = Class.PDI + "%"
    replace_text({'<PDI>':  pdi}, shapes)

    replace_text({'<ADCF>': Class.ADCF_affichage}, shapes) 

    #on rjaoute les % sans que ca ait un impact sur les calculs
    bac = Class.BAC + "%"
    bac = bac.replace(".", ",")

    replace_text({'<BAC>':  bac}, shapes)
    
    if (Class.F0 == "jours"):
        Class.BCPN = Class.BCPN + "% environ"
    else:
        Class.BCPN = Class.BCPN


    replace_text({'<BCPN>':  Class.BCPN}, shapes)

    com = str(Class.COM) + "%"
    com = com.replace(".", ",")
    replace_text({'<COM>':  com}, shapes)
    nsd = Class.NSD + "%"
    replace_text({'<NSD>':  nsd}, shapes)


    replace_text({'<first>':  Class.first}, shapes)
    replace_text({'<P2>':  Class.p2}, shapes)
    replace_text({'<last-2>':  Class.last_2}, shapes)
    replace_text({'<last-1>':  Class.last_1}, shapes)
    replace_text({'<last>':  Class.last}, shapes)

    nsm = Class.NSM + "%"
    replace_text({'<NSM>':  nsm}, shapes)
    nsf = str(Class.NSF) + "%"
    replace_text({'<NSF>':  nsf}, shapes)

    abdac = float(Class.ABDAC)
    len_cpn = (str(abdac)[::-1].find('.'))

    if (len_cpn < 2):
            cpn = (f'{float(Class.ABDAC):.2f}')

            
    abdac = str(Class.ABDAC).replace(".", ",")
    replace_text({'<ABDAC>':  abdac}, shapes)
    
    dbac = str(Class.DBAC) + "%"
    dbac = dbac.replace(".", ",")

    replace_text({'<DBAC>':  dbac}, shapes)
    
    deg = str(Class.DEG)
    deg = deg.replace(".", ",")
    replace_text({'<DEG>':  deg}, shapes)

def elementsToReplaceCalcul(Class, shapes):
    replace_text({'<balisedeg1>': Class.balisedeg1}, shapes)

    replace_text({'<balisedeg2>': Class.balisedeg2}, shapes)     
    replace_text({'<balisedeg3>': Class.balisedeg3}, shapes)
    replace_text({'<balisedeg4>': Class.balisedeg4}, shapes) 

    replace_text({'<baliseCM>': Class.baliseCM}, shapes) 
    replace_text({'<baliseCM2>': Class.baliseCM2}, shapes)
    replace_text({'<baliseCM22>': Class.baliseCM22}, shapes) 

    Class.baliseCM3 = Class.baliseCM3.replace("(1)", '⁽¹⁾')
    Class.baliseCM4 = Class.baliseCM4.replace("(1)", '⁽¹⁾')

    replace_text({'<baliseCM3>': Class.baliseCM3}, shapes) 
    replace_text({'<baliseCM4>': Class.baliseCM4}, shapes)
    replace_text({'<baliseCM5>': Class.baliseCM5}, shapes)
    replace_text({'<baliseCM6>': Class.baliseCM6}, shapes)

    replace_text({'<SV>': Class.SV}, shapes) 

    if (Class.F0 == "jours"):
        Class.DDCI = str(Class.DDCI) + " (exclu)"
    replace_text({'<DDCI>': Class.DDCI}, shapes) 

    replace_text({'<BLOCDIVIDENDE>': Class.BLOCDIVIDENDE}, shapes)
    bloc2 = str(Class.BLOCDIVIDENDE).capitalize()
    replace_text({'<BLOCDIVIDENDE2>': bloc2}, shapes)
    replace_text({'<BLOCDIVIDENDE3>': str(Class.BLOCDIVIDENDE3)}, shapes)

    deg = float(Class.DEG)
    deg = ("{:.2f}".format(deg))
        
    for ticker in Class.Yahoo_value_name:
        Class.legende_tickers =  Class.legende_tickers + "\n" + ticker

    replace_text({'<WALLY>': str(Class.WALLY)}, shapes)
    replace_text({'<SFPASANNEE>': str(Class.SFPASANNEE)}, shapes)

    replace_text({'<DUREE>': str(Class.DUREE)}, shapes)
    replace_text({'<ENTRE>': str(Class.entre)}, shapes)

    replace_text({'<BRA>': str(Class.BRA)}, shapes)
    replace_text({'<BVC>': str(Class.BVC)}, shapes)
    replace_text({'<ANNUALISE>': str(Class.ANNUALISE)}, shapes)

    replace_text({'<F0s>':  Class.F0s}, shapes)


    replace_text({'<1PDC>': Class.PDC1_affichage}, shapes) 
    
    replace_text({'<2PDC>': Class.PDC2_affichage}, shapes) 
    replace_text({'<DDR>': Class.DDR}, shapes)

    ddr1 = str(Class.DDR1)[0:10]
    ddr1 = ddr1[8:10] + "/"+ ddr1[5:7] + "/" + ddr1[0:4]
    replace_text({'<DDR1>': ddr1}, shapes)
    replace_text({'<DDR1_MAJ>': Class.DDR1_MAJ}, shapes)
    replace_text({'<DDR1_MAJ_MIN>': Class.DDR1_MAJ_MIN}, shapes)

    replace_text({'<DDR1-12_MAJ>': Class.DDR1_12_MAJ}, shapes)
    replace_text({'<DDR1-12>': Class.DDR1_12}, shapes)

    replace_text({'<DIC>': Class.DIC}, shapes)
    
    pdiperf = str(int(Class.PDIPERF)) + "%"
    pdiperf = pdiperf.replace(".", ",")


    replace_text({'<PDIPERF>': pdiperf}, shapes)
    replace_text({'<P>': Class.prefix}, shapes)

    pr1= str(Class.PR1)
    replace_text({'<1PR>': str(Class.PR1)},  shapes)
    replace_text({'<DPRR>': Class.DPRR}, shapes)
    replace_text({'<TDP>': Class.TDP}, shapes)
    replace_text({'<GC>': Class.GC}, shapes)
    gca = str(Class.GCA) + "%"
    gca = gca.replace(".", ",")

    replace_text({'<GCA>': gca}, shapes)
    gce = ("{:.2f}".format(Class.GCE))
    gce = str(gce) + "%"
    gce = gce.replace(".", ",")

    if (str(Class.inconvenient == "NULL")):
        Class.inconvenient = ""
 
    replace_text({'<inconv>': Class.inconvenient}, shapes)
    replace_text({'<inconvénient>': Class.inconvenient}, shapes)
    replace_text({'<R1>': Class.R1}, shapes)
    replace_text({'<GCE>': gce}, shapes)
    replace_text({'<ABAC>': Class.ABAC}, shapes)
    replace_text({'<decrement>': Class.decrement}, shapes)

    replace_text({'<NDR>': Class.NDR }, shapes)
    replace_text({'<ADPR>': str(Class.ADPR)}, shapes)
    replace_text({'<SJR1>': Class.SJR1}, shapes)
    replace_text({'<SJR2>': Class.SJR2}, shapes)
    replace_text({'<SJR3>': Class.SJR3}, shapes)
    replace_text({'<SJR4>': Class.SJR4}, shapes)
    replace_text({'<SJR5>': Class.SJR5}, shapes)
    replace_text({'<SJR6>': Class.SJR6}, shapes)
    replace_text({'<SJR8>': Class.SJR8}, shapes)
    replace_text({'<SJR7>': Class.SJR7}, shapes)

    replace_text({'<SJR1M>': Class.SJR1M}, shapes)
    replace_text({'<SJR2M>': Class.SJR2M}, shapes)
    replace_text({'<SJR3M>': Class.SJR3M}, shapes)
    replace_text({'<SJR4M>': Class.SJR4M}, shapes)
    replace_text({'<SJR5M>': Class.SJR5M}, shapes)
    replace_text({'<SJR6M>': Class.SJR6M}, shapes)
    replace_text({'<SJR7M>': Class.SJR7M}, shapes)
    replace_text({'<SJR8M>': Class.SJR8M}, shapes)
    
    replace_text({'<F1>': Class.F1}, shapes)
    replace_text({'<F1_MAJ>':  Class.F1_MAJ}, shapes)

    replace_text({'<F2>': Class.F2}, shapes)

    replace_text({'<TDS>': Class.TDS}, shapes)
    replace_text({'<DU>': Class.DU}, shapes)
    replace_text({'<DU1>':  Class.DU1}, shapes)

    replace_text({'<sponsor>': Class.sponsor}, shapes)
    replace_text({'<SPONSOR>': Class.SPONSOR}, shapes)

    replace_text({'<TICKER>': Class.TICKER}, shapes)
    replace_text({'<NOMSOUSJACENT>': Class.NOMSOUSJACENT}, shapes)
    replace_text({'<DIVIDENDE>': Class.DIVIDENDE}, shapes)
    replace_text({'<SITE>': Class.Site}, shapes)
    replace_text({'<DPCI>': Class.DPCI}, shapes) 
    replace_text({'<DPCI_MAJ>':  Class.DPCI_MAJ}, shapes)
    replace_text({'<environ>':  Class.environ}, shapes)
    replace_text({'<exclus>':  Class.exclus}, shapes)

    replace_text({'<NOMP1>': Class.NOMP1}, shapes) 
    replace_text({'<NOMSOUSJACENTP1>': Class.NOMSOUSJACENTP1}, shapes) 
    replace_text({'<SJR6P1>': Class.SJR6P1}, shapes) 
    replace_text({'<type_bar2>': Class.type_bar2}, shapes) 
    replace_text({'<1PR-1>': Class.PR1_1}, shapes) 
    replace_text({'<ABAC2>': Class.ABAC2}, shapes) 
    replace_text({'<ABAC2_MAJ>': Class.ABAC2_MAJ}, shapes)
    replace_text({'<gainoucoupon>': Class.GainOuCoupon}, shapes) 
    replace_text({'<gain>': Class.GainOuCoupon}, shapes) 
    replace_text({'<PERIODE_DE_REMBOURSEMENT>': Class.PERIODE_DE_REMBOURSEMENT}, shapes) 
    replace_text({'<PERIODE_DE_REMBOURSEMENT2>': Class.PERIODE_DE_REMBOURSEMENT2}, shapes) 
    replace_text({'<PERIODE_DE_REMBOURSEMENT_MAJ>': Class.PERIODE_DE_REMBOURSEMENTM}, shapes) 

    replace_text({'<DDP>': Class.DDP}, shapes) 
    replace_text({'<Mémoire>': Class.Memoire}, shapes) 
    replace_text({'<Mémoire2>': Class.Memoire2}, shapes) 
    replace_text({'<Mémoire3>': Class.Memoire3}, shapes) 
    replace_text({'<Mémoire4>': Class.Memoire4}, shapes)
    replace_text({'<Mémoire5>': Class.Memoire5}, shapes) 
    replace_text({'<Mémoire6>': Class.Memoire6}, shapes) 
    bfp = str(Class.BFP) + "%"
    replace_text({'<BFP>': bfp}, shapes) 
    replace_text({'<PAGE>': Class.PAGE}, shapes) 
    replace_text({'<DIVERSACTION>': Class.DIVERSACTION}, shapes) 

    replace_text({'<DDPP>': Class.DDPP}, shapes) 
    replace_text({'<legendeticker>':  Class.legende_tickers}, shapes)
    replace_text({'<ticker>':  Class.legende_tickers}, shapes)
    replace_text({'<worst>':  Class.worst}, shapes)

    cpr1 = str(Class.CPR1) + "%"
    cpr1 = cpr1.replace(".", ",")
    replace_text({'<CPR1>': cpr1}, shapes)
    replace_text({'<ADDPLUSIFAUTOCALL>': Class.ADDPLUSIFAUTOCALL}, shapes) 


    #CHANGER LE NOM DE LA BALISE ET DE LA CLASS dans myclass.py
    replace_text({'<balise>': Class.balise}, shapes)
    TRA_replace(Class, shapes)

def hardcode_replace(Class, shapes):
    replace_text({"l' année": "l'année"}, shapes)
    replace_text({" , ": ", "}, shapes)
    replace_text({"(1)": '\u00281\u0029'}, shapes)
    replace_text({"%  ": '% '}, shapes)
    replace_text({"%   ": '% '}, shapes)
    replace_text({"% %": '% '}, shapes)
    replace_text({"  ": ' '}, shapes)


    replace_text({".\u00281\u0029.": '\u00281\u0029.'}, shapes)

    replace_text({"(2)": '\u00282\u0029'}, shapes)
    replace_text({"%%%": '%'}, shapes)
    replace_text({"%%": '%'}, shapes)

    replace_text({" .": '.'}, shapes)

def TRA_replace(Class, shapes):
    tra_a_s1 = Class.TRA_A_S1.replace(".", ",")
    tra_a_s1 = tra_a_s1 + "%"
    replace_text({'<TRA.D.A>': tra_a_s1}, shapes)
    tra_a_s2_100 = Class.TRA_A_S2_100.replace(".", ",")
    tra_a_s2_100 = tra_a_s2_100 + "%"
    replace_text({'<TRA.M.A>': tra_a_s2_100}, shapes)

    tra_a_s2_gain = Class.TRA_A_S2_GAIN.replace(".", ",")
    tra_a_s2_gain = tra_a_s2_gain + "%"
    replace_text({'<TRA.MG.A>': tra_a_s2_gain}, shapes)

    tra_a_m_sj = Class.TRA_M_SJ.replace(".", ",")
    tra_a_m_sj = tra_a_m_sj + "%"
    replace_text({'<TRA.M.SJ>': tra_a_m_sj}, shapes)

    tra_a_f_a = Class.TRA_F_A.replace(".", ",")
    tra_a_f_a = tra_a_f_a + "%"
    replace_text({'<TRA.F.A>': tra_a_f_a}, shapes)

    tra_f_sj = Class.TRA_F_SJ.replace(".", ",")
    tra_f_sj = tra_f_sj + "%"
    replace_text({'<TRA.F.SJ>': tra_f_sj}, shapes)

    tra_min_a = Class.TRA_MIN_A.replace(".", ",")
    tra_min_a = tra_min_a + "%"
    replace_text({'<TRA.MRA.MIN.A>': tra_min_a}, shapes)

    tra_echeance_perte_a = Class.TRA_echeance_perte_A.replace(".", ",")
    tra_echeance_perte_a = tra_echeance_perte_a + "%"
    replace_text({'<TRA.ECHEANCE.PERTE.A>': tra_echeance_perte_a}, shapes)

    # cmtra  = Class.BaliseCMTRA.replace(".", ",")
    cmtra =  Class.BaliseCMTRA

    replace_text({'<BALISECMTRA>': cmtra}, shapes)
    
    tra_tout_p = Class.TRA_TOUT_P.replace(".", ",")
    tra_tout_p = tra_tout_p + "%"
    tra_d_p = Class.TRA_D_P.replace(".", ",")
    tra_d_p = tra_d_p + "%"   
    replace_text({'<TRA.D.P>': tra_d_p}, shapes)

    tra_m_p = Class.TRA_M_P.replace(".", ",")
    tra_m_p = tra_m_p + "%"   
    replace_text({'<TRA.M.P>': tra_m_p}, shapes)

  
    tra_m_pm = Class.TRA_M_PM.replace(".", ",")
    tra_m_pm = tra_m_pm + "%"     
    replace_text({'<TRA.M.PM>': tra_m_pm}, shapes)

    tra_gm_p = Class.TRA_GM_P.replace(".", ",")
    tra_gm_p = tra_gm_p + "%"      
    replace_text({'<TRA.GM.P>': tra_gm_p}, shapes)

    tra_gm_pm = Class.TRA_GM_PM.replace(".", ",")
    tra_gm_pm = tra_gm_pm + "%"    
    replace_text({'<TRA.GM.PM>': tra_gm_pm}, shapes)
    
    tra_mra_min_p = Class.TRA_MRA_MIN_P.replace(".", ",")
    tra_mra_min_p = tra_mra_min_p + "%"    
    replace_text({'<TRA.MRA.MIN.P>': tra_mra_min_p}, shapes)
    
    tra_f_p = str(Class.TRA_F_P).replace(".", ",")
    tra_f_p = tra_f_p + "%"   
    replace_text({'<TRA.F.P>': tra_f_p}, shapes)

    tra_min_pm = Class.TRA_MRA_MIN_PM.replace(".", ",")
    tra_min_pm = tra_min_pm + "%"   
    replace_text({'<TRA.MRA.MIN.PM>': tra_min_pm}, shapes)

    tra_tout_1_p = Class.TRA_TOUT_1_P.replace(".", ",")
    tra_tout_1_p = tra_tout_1_p + "%"
    replace_text({'<TRA.TOUT-1.P>': tra_tout_1_p}, shapes)

    tra_mra_p = Class.TRA_MRA_P.replace(".", ",")
    tra_mra_p = tra_mra_p + "%" 
    replace_text({'<TRA.MRA.MAX.P>': tra_mra_p}, shapes)

    tra_med_p = Class.TRA_MED_P.replace(".", ",")
    tra_med_p = tra_med_p + "%" 
    replace_text({'<TRA.MED.P>': tra_med_p}, shapes)
    tra_tout_sauf_p = Class.TRA_TOUT_SAUF_P.replace(".", ",")
    tra_tout_sauf_p = tra_tout_sauf_p + "%" 
    replace_text({'<TRA.TOUT.SAUF.P>': tra_tout_sauf_p}, shapes)
    tra_tout_p = Class.TRA_TOUT_P.replace(".", ",")
    tra_tout_p = tra_tout_p + "%" 
    replace_text({'<TRA.TOUT.P>': tra_tout_p}, shapes)
    tra_mre_min_p = Class.TRA_MRE_MIN_P.replace(".", ",")
    tra_mre_min_p = tra_mre_min_p + "%" 
    replace_text({'<TRA.MRE.MIN.P>': tra_mre_min_p}, shapes)
    tra_mre_min_pm = Class.TRA_MRE_MIN_PM.replace(".", ",")
    tra_mre_min_pm = tra_mre_min_pm + "%" 
    replace_text({'<TRA.MRE.MIN.PM>': tra_mre_min_pm}, shapes)
    tra_max_p = str(Class.TRA_MAX_P).replace(".", ",")
    tra_max_p = tra_max_p + "%" 
    replace_text({'<TRA.MAX.P>':tra_max_p}, shapes)
    tra_em_p = Class.TRA_EM_P.replace(".", ",")
    tra_em_p = tra_em_p + "%" 
    replace_text({'<TRA.EM.P>': tra_em_p}, shapes)
    tra_rm_p = Class.TRA_RM_P.replace(".", ",")
    tra_rm_p = tra_rm_p + "%" 
    replace_text({'<TRA.RM.P>': tra_rm_p}, shapes)

    tra_mp = Class.TRA_MP.replace(".", ",")
    tra_mp = tra_mp + "%" 
    replace_text({'<TRA.MP>': tra_mp}, shapes)

    Dates_maj(Class, shapes)

def Dates_maj(Class, shapes):
    replace_text({'<1PDC_MAJ>': Class.PDC1_MAJ}, shapes)
    replace_text({'<2PDC_MAJ>': Class.PDC2_MAJ}, shapes)
    replace_text({'<DDR_MAJ>': Class.DDR_MAJ}, shapes)
    replace_text({'<DPR_MAJ>': Class.DPR_MAJ}, shapes)

    replace_text({'<DEC_MAJ>': Class.DEC_MAJ}, shapes)
    Replace_Boucle_Dates(Class, shapes)

def Replace_Boucle_Dates(Class, shapes):
    #Si frequence = jours on change la liste (deja en string) de dates en phrases
    if Class.F0 == "jours": 
        Class.dates_constat_autocall = "Chaque jour de bourse entre le " + Class.DPR_MAJ + " (inclus) et le " + Class.DCF_MAJ +"."
        Class.dates_paiement_autocall = "Le " + str(Class.NJO) + "e jour ouvré suivant la date de constatation quotidienne."
        Class.Datesconstatations1 = Class.dates_constat_autocall
        Class.Datesremb1 = Class.dates_paiement_autocall
        date_constatation =  Class.Datesconstatations1
        datesremb1 = Class.Datesremb1
        date_constatation3 =  "Chaque jour de bourse entre le " + Class.DPR_MAJ + " (inclus) et le " + Class.DCF_MAJ +"."
        Class.Datespaiement1 = "Chaque jour de bourse entre le " + Class.DPR_MAJ + " (inclus) et le " + Class.DCF_MAJ +"."
    

    #Creation d une variable de date (en string) temporaire egale a DDCI qu on convertit en format francais
    tmp = Class.DDCI #derniere date de constat initiale
    tmp = tmp[8:10] + "/" + tmp[5:7] + "/" + tmp[0:4]

    #On jarte la premiere date de constat si la date de constatation est la même que DDCI
    if (str(Class.Datesconstatations1[0:11]) == tmp):  
                Class.Datesconstatations1 = (Class.Datesconstatations1[11:])    

    #Si la frequence est en mois, alors comme en jours on change la liste de dates en string
    elif (Class.F0 == "mois"):
        date_premier =  Class.dates_paiement_phoenix[0:11]
        jours = Class.DEC[-2:]

        jour_constat = Class.DDCI[-2:]
        date_premier_phoenix_constat = Class.dates_constat_phoenix[0:11]
        date_premier_autocall_constat = Class.dates_constat_autocall[0:11]

        #si c est un phoenix ou athena, les phrases diffèrent
        if (Class.Typologie == "coupon phoenix"):
            Class.dates_constat_phoenix = "Tous  les " + jour_constat + " de chaque mois à partir du " + str(date_premier_phoenix_constat) + " (inclus), jusqu'au " + Class.DCF_affichage + " (inclus), ou le Jour de Bourse suivant si le " + jour_constat + " du mois n'est pas un Jour de Bourse. Toutes les dates peuvent faire l'objet d'ajustements en cas de jours non ouvrés."
            date_constatation = "Tous  les " + jour_constat + " de chaque mois à partir du " + str(date_premier_phoenix_constat) + " (inclus), jusqu'au " + Class.DCF_affichage + " (inclus), ou le Jour de Bourse suivant si le " + jour_constat + " du mois n'est pas un Jour de Bourse. Toutes les dates peuvent faire l'objet d'ajustements en cas de jours non ouvrés."
            date_constatation3 = date_constatation
            Class.Datespaiement1 = "Le " + str(Class.NJO) + "e Jour de Bourse suivant la date de constation mensuelle. Toutes les dates peuvent faire l'objet d'ajustements en cas de jours non ouvrés."
        else:
            Class.dates_constat_autocall = "Tous  les " + jour_constat + " de chaque mois à partir du " + str(date_premier_autocall_constat) + " (inclus), jusqu'au " + Class.DCF_affichage + " (inclus), ou le Jour de Bourse suivant si le " + jour_constat + " du mois n'est pas un Jour de Bourse. Toutes les dates peuvent faire l'objet d'ajustements en cas de jours non ouvrés."
            date_constatation = "Tous  les " + jour_constat + " de chaque mois à partir du " + str(date_premier_autocall_constat) + " (inclus), jusqu'au " + Class.DCF_affichage + " (inclus), ou le Jour de Bourse suivant si le " + jour_constat + " du mois n'est pas un Jour de Bourse. Toutes les dates peuvent faire l'objet d'ajustements en cas de jours non ouvrés."
            date_constatation3 = date_constatation
            Class.Datespaiement1 = "Le " + str(Class.NJO) + "e Jour de Bourse suivant la date de constation mensuelle. Toutes les dates peuvent faire l'objet d'ajustements en cas de jours non ouvrés."  
    
    
    #Si la frequence n est ni en jours ni en mois
    else:
        date_constatation = Class.Datesconstatations1
        date_constatation3 = Class.Datesconstatations3
    
    #si la frequence est en mois, on continue de changer la liste de dates en phrases
    if (Class.F0 == "mois"):
        jours = Class.DDCI[-2:]
        datesremb1 = "Le " + str(Class.NJO) + "e Jour de Bourse suivant la date de constation mensuelle à partir du " + Class.DR1_affichage + " (inclus). Toutes les dates peuvent faire l'objet d'ajustements en cas de jours non ouvrés."
        datesremb3 = datesremb1
    else:
        datesremb1 = Class.Datesremb1
        datesremb3 = datesremb1

    replace_text({'<Datesconstatations1>': date_constatation}, shapes)
    replace_text({'<Datesconstatations2>': Class.Datesconstatations2}, shapes)

    replace_text({'<Datesconstatations3>': date_constatation3}, shapes)
    replace_text({'<Datesconstatations4>': Class.Datesconstatations4}, shapes)

    replace_text({'<Datesremb1>': datesremb1}, shapes)
    replace_text({'<Datesremb2>': Class.Datesremb2}, shapes)
    replace_text({'<Datesremb3>': datesremb3}, shapes)
    replace_text({'<Datesremb4>': Class.Datesremb4}, shapes)
    replace_text({'<Datesremb5>': Class.Datesremb5}, shapes)
    replace_text({'<Datesremb6>': Class.Datesremb6}, shapes)
    replace_text({'<Datesremb7>': Class.Datesremb7}, shapes)
    replace_text({'<Datesremb8>': Class.Datesremb8}, shapes)

    replace_text({'<Datespaiement1>': Class.Datespaiement1}, shapes)
    replace_text({'<Datespaiement2>': Class.Datespaiement2}, shapes)
    replace_text({'<Datespaiement3>': Class.Datespaiement3}, shapes)
    replace_text({'<Datespaiement4>': Class.Datespaiement4}, shapes)
    replace_text({'<Datespaiement5>': Class.Datespaiement5}, shapes)
    replace_text({'<Datespaiement6>': Class.Datespaiement6}, shapes)
    replace_text({'<Datespaiement7>': Class.Datespaiement7}, shapes)
    replace_text({'<Datespaiement8>': Class.Datespaiement8}, shapes)

    replace_text({'<dates_constat_autocall>': Class.dates_constat_autocall}, shapes)
    replace_text({'<dates_paiement_autocall>': Class.dates_paiement_autocall}, shapes)
    replace_text({'<dates_constat_phoenix>': Class.dates_constat_phoenix}, shapes)
    replace_text({'<dates_paiement_phoenix>': Class.dates_paiement_phoenix}, shapes)
    replace_text({'<dates_last_remboursement_rappel>': Class.dates_last_remboursement_rappel}, shapes)
    hardcode_replace(Class, shapes)

#PREMIER RAPPEL A DATE DERNIER RAPPEL

#avoir toutes les slides
def getAllSlides(prs):
    slides = [slide for slide in prs.slides]
    shapes = []

    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    return shapes



def add_row(table) -> _Row:
    new_row = deepcopy(table._tbl.tr_lst[-1]) 
    # duplicating last row of the table as a new row to be added

    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = '' # defaulting cell contents to empty text

    table._tbl.append(new_row) 
    return table

    
#remplace le texte sur le ppt 
def replace_text(replacements: dict, shapes: List):
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:

                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text


            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if match in cell.text:
                            new_text = cell.text.replace(match, replacement)
                            cell.text = new_text


#load le ppt, remplace les balises et le sauvegarde
def ChangeTextOnPpt(Class):
    #NAME = "result/"+ "BC - " + Class.Nom + " - " + Class.Isin + ".pptx" 
    NAME = f'''result/BC - {Class.Nom} - {Class.Isin}.pptx'''
    prs_string = "templates/"+ Class.template + ".pptx" 
    prs = Presentation(prs_string)
    shapes = getAllSlides(prs)
    Class.shapes = shapes

    elementsToReplaceDegressivite(Class, shapes)

    elementsToReplaceRemplacement(Class, shapes)
    elementsToReplaceCalcul(Class, shapes)

    #on repasse une 2 eme fois au cas ou certaines transformations soient mal placées
    elementsToReplaceDegressivite(Class, shapes)

    elementsToReplaceRemplacement(Class, shapes)
    elementsToReplaceCalcul(Class, shapes)

    compteur = 0
    for shape in Class.shapes:
        textbox = shape
        for i in Class.deleteblocs:

            try:
                if (i in textbox.text):
                        print("paragraphe supprimé")
                        compteur+=1
                        sp = textbox.element
                        sp.getparent().remove(sp)

            except Exception:
                pass
            
    if compteur > 1:
        print(compteur, "paragraphes ont été supprimés")
    else:

        print(compteur, "paragraphe a été supprimé")
    # from colormath.color_objects import sRGBColor
    compteur = 0
    compteur_tab = 1

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                
                table = shape.table
            
                numberoftickers = len(Class.Yahoo) -1 #pour savoir le nombre de rows a ajouter
                if(compteur_tab == 1 and len(Class.Yahoo) > 1):  #On ajoute des rows pour le tableau de performance des tickers yahoo
                    for newcolumn in range((numberoftickers)):
                        table = add_row(table)

                for row_idx, row in enumerate(table.rows): 
                    for col_idx, cell in enumerate(row.cells):
                        cell = table.cell(row_idx, col_idx)
                        # cell.text = 'TEST DE LA STREET'
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(7)

                        if (compteur_tab == 1 and col_idx == 0 and row_idx > 0): #on ajoute les sousjacent a la premiere colonne de chaque ligne
                            try:

                                cell = table.cell(row_idx, col_idx)
                            
                                # cell.text = str(Class.Yahoo_value_name[row_idx -1]) +" " + str(Class.Yahoo_value_dividende[row_idx -1])
                                cell_tmp = Class.NOMSOUSJACENTP1.split('ET')
                                cell.text = cell_tmp[row_idx - 1]
                                paragraph.font.size = Pt(7)

                            except Exception:
                                print("pas de tableau de perf")

                        # print("%r is cells[%d][%d]" % (cell, row_idx, col_idx))
                        cell = table.cell(row_idx, col_idx)
                        # cell.text = 'TEST DE LA STREET'
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(7)
                        # paragraph.font.color = 'rgb(0, 0, 0)'

                        # paragraph.font.color.rgb = screen.bgcolor(color)

                        if(compteur_tab ==1 and row_idx >= 1 and col_idx >= 1 ):
                            try:
                                cell.text = Class.Yahoo_value[row_idx - 1][col_idx - 1]
                                paragraph = cell.text_frame.paragraphs[0]
                                paragraph.font.size = Pt(7)
                            except Exception: 
                                print("pas de text sur les paragraphes")
                compteur_tab +=1


            if shape.has_text_frame:
                if ("<graph1>" in shape.text):
                    cur_text = shape.text
                    new_text = cur_text.replace(str("<graph1>"), str(""))
                    shape.text = new_text
                    pic = slide.shapes.add_picture("graph1.png", Inches(0), Inches(6.71), Inches(8))
                
                if ("<graph2>" in shape.text):  
                    cur_text = shape.text
                    new_text = cur_text.replace(str("<graph2>"), str(""))
                    shape.text = new_text
                    pic = slide.shapes.add_picture("graph_scenario_def.png", Inches(0.45), Inches(1.8), Inches(3.7))

                if ("<graph3>" in shape.text):  
                    cur_text = shape.text
                    new_text = cur_text.replace(str("<graph3>"), str(""))
                    shape.text = new_text
                    pic = slide.shapes.add_picture("graph_scenario_median.png", Inches(0.45), Inches(4.78), Inches(3.6))
    
                if ("<graph4>" in shape.text):  
                    cur_text = shape.text
                    new_text = cur_text.replace(str("<graph4>"), str(""))
                    shape.text = new_text
                    pic = slide.shapes.add_picture("graph_scenario_fav.png", Inches(0.45), Inches(7.68), Inches(3.45))
                   
                if ("<graph5>" in shape.text):  
                        cur_text = shape.text
                        new_text = cur_text.replace(str("<graph5>"), str(""))
                        shape.text = new_text
                        pic = slide.shapes.add_picture("graph5.png", Inches(0.35), Inches(4.7), width=Inches(7.5), height=Inches(4.3))
                        
                if ("graph_legende.png" in shape.text):
                    cur_text = shape.text
                    new_text = cur_text.replace(str("graph_legende.png"), str(""))
                    shape.text = new_text
                    pic = slide.shapes.add_picture("graph_legende.png", Inches(0), Inches(6.71), Inches(8))

    try:
        print("Nettoyage du projet, supression des documents inutiles")
        os.remove("graph1.png")
        print("Suppression de graph1.png")
        os.remove("graph2.png")
        print("Suppression de graph2.png")
        os.remove("graph3.png")
        print("Suppression de graph3.png")
        os.remove("graph4.png")
        print("Suppression de graph4.png")
        os.remove("graph5.png")
        print("Suppression de graph5.png")

    except Exception:
        print("OH NONNNN")

    # xml_slides = presentation.slides._sldIdLst  
    #slides = list(prs)
    # xml_slides.remove(slides[index])
    # 

    excel(Class) 
    # prs.slides.remove(0)
   
    if(Class.Typologie == "coupon autocall"):  #SUPPRIMER LES PAGES EN FONCTION DE SI C EST ATHENA OU PHOENIX athena: 3,5,6,8;10, 13
        rId = prs.slides._sldIdLst[12].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[12]

        rId = prs.slides._sldIdLst[9].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[9]

        rId = prs.slides._sldIdLst[7].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[7]
    
        rId = prs.slides._sldIdLst[5].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[5]

        rId = prs.slides._sldIdLst[4].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[4]
        
        rId = prs.slides._sldIdLst[2].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[2]
    
    else:
        rId = prs.slides._sldIdLst[11].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[11]

        rId = prs.slides._sldIdLst[8].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[8]

        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]
        
        rId = prs.slides._sldIdLst[3].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[3]

        rId = prs.slides._sldIdLst[1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[1]
    hardcode_replace(Class, shapes)
    prs.save(NAME)


    #le ppt de sauvegarde
    NAME = "result/Base Graphique "+ Class.Nom + " - " + Class.Isin + ".pptx" 
    prs_string = "templates/Base Graphique.pptx" 
    prs = Presentation(prs_string)
    shapes = getAllSlides(prs)
    Class.shapes = shapes

    elementsToReplaceDegressivite(Class, shapes)

    elementsToReplaceRemplacement(Class, shapes)
    elementsToReplaceCalcul(Class, shapes)

    #on repasse une 2 eme fois au cas ou certaines transformations soient mal placées
    elementsToReplaceDegressivite(Class, shapes)

    elementsToReplaceRemplacement(Class, shapes)
    elementsToReplaceCalcul(Class, shapes)

    prs.save(NAME)

